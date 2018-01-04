#!/usr/bin/env python

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "graph test"

graph_area = slide.shapes.placeholders[1]
graph_area.text = "graph test"
x0 = graph_area.left
y0 = graph_area.top
xs = graph_area.width
ys = graph_area.height

def set_style(shape, name):
    shape.text_frame.paragraphs[0].text = name
    font = shape.text_frame.paragraphs[0].font
    font.size = Pt(12)
    font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    shape.line.width = Pt(1)
    #shape.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    shape.fill.background()

def connect(shape0, shape1):
    edge = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, 0, 0, 0, 0)
    edge.begin_connect(shape0, 3)
    edge.end_connect(shape1, 1)
    #edge.color.rgb = RGBColor(0x00, 0x00, 0x00)

node0 = slide.shapes.add_shape(
    MSO_SHAPE.OVAL, xs/10*0 + x0, ys/10*2 + y0, xs/10*1, ys/10*1
)
set_style(node0, "node0")

label = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, xs/10*2 + x0, ys/10*4 + y0, xs/10*1, ys/10*1
)
set_style(label, "label")

node1 = slide.shapes.add_shape(
    MSO_SHAPE.OVAL, xs/10*4 + x0, ys/10*2 + y0, xs/10*1, ys/10*1
)
set_style(node1, "node1")

connect(node0, label)
connect(label, node1)

prs.save('test.pptx')
